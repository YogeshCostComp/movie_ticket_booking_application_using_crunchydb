"""
Generate a detailed comparison PPT: OpenClaw vs SRE Ephemeral Agent Architecture
4-column comparison format with advantages of each approach
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────
BG_DARK       = RGBColor(0x0F, 0x11, 0x1A)
BG_CARD       = RGBColor(0x1A, 0x1D, 0x2E)
BG_CARD_ALT   = RGBColor(0x14, 0x17, 0x25)
ACCENT_CYAN   = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_GREEN  = RGBColor(0x00, 0xE6, 0x76)
ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC)
ACCENT_RED    = RGBColor(0xFF, 0x45, 0x45)
ACCENT_YELLOW = RGBColor(0xFF, 0xD6, 0x00)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xB0, 0xB8, 0xC8)
MID_GRAY      = RGBColor(0x6B, 0x72, 0x80)
OPENCLAW_CLR  = RGBColor(0xFF, 0x45, 0x00)  # Lobster red-orange
SRE_CLR       = RGBColor(0x00, 0xD4, 0xFF)  # Cyan

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    # Radius
    if radius:
        shape.adjustments[0] = radius
    else:
        shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=13, default_color=LIGHT_GRAY):
    """lines: list of (text, size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)
    return txBox


# ── Build the 4-column comparison table on a slide ──────────────
def add_comparison_table(slide, title, rows, left_start, top_start, col_widths, row_height=Inches(0.55)):
    """
    rows: list of (aspect, openclaw_val, sre_val, verdict)
    col_widths: [aspect_w, openclaw_w, sre_w, verdict_w]
    """
    headers = ["Aspect", "🦞 OpenClaw", "⚡ SRE Ephemeral Agent", "Winner / Notes"]
    header_colors = [MID_GRAY, OPENCLAW_CLR, SRE_CLR, ACCENT_PURPLE]

    y = top_start
    # Header row
    x = left_start
    for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
        shape = add_shape(slide, x, y, cw, Inches(0.45), BG_CARD, header_colors[i])
        shape.text_frame.word_wrap = True
        p = shape.text_frame.paragraphs[0]
        p.text = hdr
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].space_before = Pt(2)
        x += cw

    y += Inches(0.5)

    # Data rows
    for idx, (aspect, oc_val, sre_val, verdict) in enumerate(rows):
        bg = BG_CARD if idx % 2 == 0 else BG_CARD_ALT
        x = left_start
        vals = [aspect, oc_val, sre_val, verdict]
        colors = [ACCENT_YELLOW, LIGHT_GRAY, LIGHT_GRAY, ACCENT_GREEN]
        bolds = [True, False, False, False]
        sizes = [10, 10, 10, 10]

        for i, (val, cw) in enumerate(zip(vals, col_widths)):
            shape = add_shape(slide, x, y, cw, row_height, bg)
            shape.text_frame.word_wrap = True
            p = shape.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(sizes[i])
            p.font.color.rgb = colors[i]
            p.font.bold = bolds[i]
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.LEFT
            shape.text_frame.margin_left = Pt(6)
            shape.text_frame.margin_top = Pt(3)
            x += cw
        y += row_height + Inches(0.02)

    return y


def build_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Decorative top bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_CYAN)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 "🦞  OpenClaw  vs  ⚡ SRE Ephemeral Agent",
                 font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(0.8),
                 "A Detailed Architecture & Design Comparison",
                 font_size=22, color=ACCENT_CYAN, bold=False, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.6),
                 "Personal AI Assistant  vs  Domain-Specific SRE Orchestrator",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Two logos / labels
    oc_box = add_shape(slide, Inches(3), Inches(4.8), Inches(3), Inches(1.2), BG_CARD, OPENCLAW_CLR)
    oc_box.text_frame.word_wrap = True
    p = oc_box.text_frame.paragraphs[0]
    p.text = "🦞 OpenClaw"
    p.font.size = Pt(22)
    p.font.color.rgb = OPENCLAW_CLR
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = oc_box.text_frame.add_paragraph()
    p2.text = "211k ★ · 692 contributors"
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    sre_box = add_shape(slide, Inches(7), Inches(4.8), Inches(3), Inches(1.2), BG_CARD, SRE_CLR)
    sre_box.text_frame.word_wrap = True
    p = sre_box.text_frame.paragraphs[0]
    p.text = "⚡ SRE Agent"
    p.font.size = Pt(22)
    p.font.color.rgb = SRE_CLR
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = sre_box.text_frame.add_paragraph()
    p2.text = "Ephemeral · MCP · Claude"
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(3), Inches(6.6), Inches(7), Inches(0.5),
                 "February 2026  ·  Yogesh Kumar  ·  IBM",
                 font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 2 — Overview: What Each System Is
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_ORANGE)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "What Are We Comparing?", font_size=28, color=WHITE, bold=True)

    # OpenClaw card
    add_shape(slide, Inches(0.4), Inches(1.2), Inches(5.9), Inches(5.5), BG_CARD, OPENCLAW_CLR)
    add_text_box(slide, Inches(0.7), Inches(1.3), Inches(5.3), Inches(0.5),
                 "🦞 OpenClaw — Personal AI Assistant", font_size=18, color=OPENCLAW_CLR, bold=True)

    oc_lines = [
        ("• Always-on personal assistant daemon", 13, LIGHT_GRAY, False),
        ("• Runs on YOUR devices (local-first)", 13, LIGHT_GRAY, False),
        ("• 10+ messaging channels (WhatsApp, Telegram,", 13, LIGHT_GRAY, False),
        ("  Slack, Discord, Signal, iMessage, Teams, etc.)", 13, LIGHT_GRAY, False),
        ("• Gateway = WebSocket control plane", 13, LIGHT_GRAY, False),
        ("• Pi Agent runtime (persistent, session-based)", 13, LIGHT_GRAY, False),
        ("• Skills/Tools ecosystem (ClawHub registry)", 13, LIGHT_GRAY, False),
        ("• Sub-agents via sessions_spawn (semi-ephemeral)", 13, ACCENT_YELLOW, True),
        ("• Voice Wake + Talk Mode + Live Canvas", 13, LIGHT_GRAY, False),
        ("• 211k stars · MIT License · TypeScript", 13, MID_GRAY, False),
    ]
    add_multi_text(slide, Inches(0.7), Inches(1.9), Inches(5.3), Inches(4.5), oc_lines)

    # SRE Agent card
    add_shape(slide, Inches(6.8), Inches(1.2), Inches(5.9), Inches(5.5), BG_CARD, SRE_CLR)
    add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.3), Inches(0.5),
                 "⚡ SRE Ephemeral Agent Orchestrator", font_size=18, color=SRE_CLR, bold=True)

    sre_lines = [
        ("• On-demand SRE specialist agents", 13, LIGHT_GRAY, False),
        ("• Runs as a web service (cloud-deployed)", 13, LIGHT_GRAY, False),
        ("• Single chat interface (WebSocket + REST)", 13, LIGHT_GRAY, False),
        ("• Orchestrator = FastAPI + LLM Brain", 13, LIGHT_GRAY, False),
        ("• 7 typed agents (Log, Health, Monitor, etc.)", 13, LIGHT_GRAY, False),
        ("• Fully ephemeral — spawn → execute → destroy", 13, ACCENT_YELLOW, True),
        ("• MCP Server integration (30+ SRE tools)", 13, LIGHT_GRAY, False),
        ("• Cooldown pool (120s) for reuse optimization", 13, ACCENT_YELLOW, True),
        ("• Real-time pipeline visualization", 13, LIGHT_GRAY, False),
        ("• Custom built · Python/FastAPI · Claude", 13, MID_GRAY, False),
    ]
    add_multi_text(slide, Inches(7.1), Inches(1.9), Inches(5.3), Inches(4.5), sre_lines)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 3 — Architecture Comparison (visual)
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "Architecture: How Each System Works", font_size=28, color=WHITE, bold=True)

    # OpenClaw architecture
    add_shape(slide, Inches(0.4), Inches(1.2), Inches(5.9), Inches(5.8), BG_CARD, OPENCLAW_CLR)
    add_text_box(slide, Inches(0.7), Inches(1.3), Inches(5.3), Inches(0.4),
                 "🦞 OpenClaw Architecture", font_size=16, color=OPENCLAW_CLR, bold=True)

    oc_arch = [
        ("WhatsApp / Telegram / Slack / Discord / ...", 11, MID_GRAY, False),
        ("              │", 11, MID_GRAY, False),
        ("              ▼", 11, MID_GRAY, False),
        ("┌─────────────────────────────┐", 11, OPENCLAW_CLR, False),
        ("│      Gateway (daemon)        │", 11, OPENCLAW_CLR, True),
        ("│   ws://127.0.0.1:18789       │", 11, OPENCLAW_CLR, False),
        ("└──────────┬──────────────────┘", 11, OPENCLAW_CLR, False),
        ("           │", 11, MID_GRAY, False),
        ("    ┌──────┼──────┐", 11, MID_GRAY, False),
        ("    ▼      ▼      ▼", 11, MID_GRAY, False),
        (" Pi Agent  Skills  Nodes", 11, WHITE, True),
        (" (persist) (tools) (devices)", 11, LIGHT_GRAY, False),
        ("    │", 11, MID_GRAY, False),
        ("    └── sessions_spawn ──► Sub-Agent", 11, ACCENT_YELLOW, True),
        ("         (semi-ephemeral, 60min archive)", 11, ACCENT_YELLOW, False),
    ]
    add_multi_text(slide, Inches(0.7), Inches(1.8), Inches(5.3), Inches(5), oc_arch)

    # SRE architecture
    add_shape(slide, Inches(6.8), Inches(1.2), Inches(5.9), Inches(5.8), BG_CARD, SRE_CLR)
    add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.3), Inches(0.4),
                 "⚡ SRE Ephemeral Architecture", font_size=16, color=SRE_CLR, bold=True)

    sre_arch = [
        ("User Chat (WebSocket)", 11, MID_GRAY, False),
        ("              │", 11, MID_GRAY, False),
        ("              ▼", 11, MID_GRAY, False),
        ("┌─────────────────────────────┐", 11, SRE_CLR, False),
        ("│   Orchestrator (FastAPI)     │", 11, SRE_CLR, True),
        ("│   + LLM Brain (Claude)       │", 11, SRE_CLR, False),
        ("└──────────┬──────────────────┘", 11, SRE_CLR, False),
        ("           │ classify intent", 11, ACCENT_YELLOW, True),
        ("    ┌──────┼──────┬──────┐", 11, MID_GRAY, False),
        ("    ▼      ▼      ▼      ▼", 11, MID_GRAY, False),
        (" Log    Health  Monitor  Runbook", 11, WHITE, True),
        (" Agent  Agent   Agent    Agent", 11, LIGHT_GRAY, False),
        ("  (spawn)  (execute)  (destroy)", 11, ACCENT_GREEN, True),
        ("    │", 11, MID_GRAY, False),
        ("    └── MCP Server (30+ SRE tools)", 11, ACCENT_PURPLE, True),
    ]
    add_multi_text(slide, Inches(7.1), Inches(1.8), Inches(5.3), Inches(5), sre_arch)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 4 — Core Architecture Comparison Table
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_CYAN)

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
                 "Core Architecture Comparison (4-Column)", font_size=26, color=WHITE, bold=True)

    col_w = [Inches(2.2), Inches(3.8), Inches(3.8), Inches(3.0)]
    rows = [
        ("Agent Lifetime",
         "Persistent daemon — always running via launchd/systemd",
         "Ephemeral — created per-query, destroyed after 120s cooldown",
         "SRE ✓ for isolation; OC ✓ for continuity"),
        ("Central Controller",
         "Gateway (WebSocket control plane on port 18789)",
         "Orchestrator (FastAPI + WebSocket)",
         "Similar concept, different scale"),
        ("Agent Runtime",
         "Single Pi Agent (embedded pi-mono) with sessions",
         "7 typed agent classes (Log, Health, Monitor, etc.)",
         "SRE ✓ specialization; OC ✓ flexibility"),
        ("Multi-Agent",
         "sessions_spawn → sub-agent (isolated session, generic)",
         "LLM classifies → typed specialist agent",
         "SRE ✓ typed routing"),
        ("Who Decides Routing",
         "The agent itself decides when to spawn helpers",
         "External LLM Brain classifies user intent",
         "Different philosophy"),
        ("Agent-to-Agent Comms",
         "Yes — sessions_send with ping-pong reply loops",
         "No — orchestrator is single coordinator",
         "OC ✓ for complex workflows"),
        ("State Management",
         "Session-persistent (JSONL transcripts, compaction)",
         "Stateless per-agent (fresh context every time)",
         "SRE ✓ for SRE; OC ✓ for chat"),
        ("Tool Ecosystem",
         "Skills platform (ClawHub), bundled/managed/workspace",
         "MCP Server with 30+ fixed SRE tools",
         "OC ✓ extensibility"),
        ("Channels / I/O",
         "10+ channels (WhatsApp, Telegram, Slack, Discord...)",
         "Single web chat + REST API",
         "OC ✓ multi-channel"),
        ("Deployment",
         "Local daemon (npm/pnpm) + Docker + Nix",
         "Docker on cloud (Render, IBM CE)",
         "Different targets"),
    ]

    add_comparison_table(slide, "Core Architecture", rows,
                         Inches(0.3), Inches(0.85), col_w, Inches(0.58))

    # ════════════════════════════════════════════════════════════════
    # SLIDE 5 — Agent Lifecycle Comparison Table
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_PURPLE)

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
                 "Agent Lifecycle & Session Comparison", font_size=26, color=WHITE, bold=True)

    rows2 = [
        ("Creation Trigger",
         "Gateway starts → Pi Agent always exists; sub-agents via sessions_spawn",
         "User query → LLM classifies intent → typed agent instantiated",
         "SRE ✓ on-demand efficiency"),
        ("Destruction",
         "Main agent: never (daemon). Sub-agents: archive after 60 min",
         "All agents: destroy after 120s cooldown (configurable)",
         "SRE ✓ aggressive cleanup"),
        ("Session Reuse",
         "Same session persists daily (reset at 4 AM or idle timeout)",
         "Cooldown pool reuses same agent within 120s window",
         "Both optimize differently"),
        ("Nesting / Depth",
         "Sub-agents cannot spawn further sub-agents",
         "Agents cannot spawn other agents — orchestrator only",
         "Same restriction, both safe"),
        ("Memory / Context",
         "Full session history (JSONL), compaction, memory flush",
         "Fresh each time — no history between queries",
         "OC ✓ continuity; SRE ✓ safety"),
        ("Registry / Audit",
         "Session store with metadata, transcripts, token counts",
         "AgentRegistry singleton with object IDs, memory, threads, PID",
         "SRE ✓ granular lifecycle proof"),
        ("Cleanup Model",
         "keep (default) or delete flag on sub-agents",
         "Always destroyed + garbage collected after cooldown",
         "SRE ✓ zero resource leakage"),
        ("Fault Isolation",
         "Sub-agent failure contained; main agent unaffected",
         "Any agent failure isolated; orchestrator routes to new one",
         "Both ✓ good isolation"),
        ("Concurrency",
         "Serialized per session key + optional global queue",
         "Multiple typed agents can run in parallel",
         "SRE ✓ parallel by design"),
        ("Cold Start Cost",
         "None — agent always running",
         "~ms (Python object creation) — mitigated by cooldown pool",
         "OC ✓ zero cold start"),
    ]

    add_comparison_table(slide, "Agent Lifecycle", rows2,
                         Inches(0.3), Inches(0.85), col_w, Inches(0.58))

    # ════════════════════════════════════════════════════════════════
    # SLIDE 6 — Security, Scalability & Ops Comparison
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_RED)

    add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
                 "Security, Scalability & Operations", font_size=26, color=WHITE, bold=True)

    rows3 = [
        ("State Leakage Risk",
         "Possible — session context carries across messages unless scoped",
         "None — every agent starts fresh, zero cross-query contamination",
         "SRE ✓ critical for incidents"),
        ("Horizontal Scaling",
         "Single gateway per install; multi-agent routing per channel",
         "Spawn N agents in parallel; no shared mutable state",
         "SRE ✓ naturally parallel"),
        ("Resource Efficiency",
         "Always-on daemon consumes base memory even when idle",
         "Zero memory when idle — agents exist only during work",
         "SRE ✓ for cloud cost"),
        ("Crash Recovery",
         "Gateway daemon restarts; sessions resume from JSONL",
         "Orchestrator stateless — just restart; no state to recover",
         "SRE ✓ simpler recovery"),
        ("Compliance / Audit",
         "Session transcripts, send policy, sandbox visibility",
         "AgentRegistry: full lifecycle audit trail per agent with timestamps",
         "Both ✓ different focus"),
        ("Sandbox Support",
         "Docker sandboxes for non-main sessions; tool allowlists",
         "Each agent is naturally sandboxed by type (tools scoped)",
         "OC ✓ explicit; SRE ✓ implicit"),
        ("DM / Access Control",
         "Pairing codes, allowlists, sendPolicy per channel/session",
         "API key auth on MCP server; single-user chat interface",
         "OC ✓ enterprise-grade access"),
        ("Monitoring",
         "openclaw doctor, health checks, usage tracking, logging",
         "Pipeline visualization, /health, AgentRegistry inspect API",
         "Both ✓ different tooling"),
        ("Updates / Hot Reload",
         "Config hot-reload (hybrid mode), skills auto-update",
         "Docker redeploy; no hot-reload needed (stateless)",
         "OC ✓ zero-downtime config"),
        ("Observability",
         "Event streams (lifecycle, assistant, tool), block streaming",
         "WebSocket pipeline events (created→thinking→executing→done)",
         "Both ✓ real-time events"),
    ]

    add_comparison_table(slide, "Security & Ops", rows3,
                         Inches(0.3), Inches(0.85), col_w, Inches(0.58))

    # ════════════════════════════════════════════════════════════════
    # SLIDE 7 — Advantages of OpenClaw
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), OPENCLAW_CLR)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "🦞 Advantages of OpenClaw", font_size=28, color=OPENCLAW_CLR, bold=True)

    advantages_oc = [
        ("Multi-Channel Reach", "10+ channels (WhatsApp, Telegram, Slack, Discord, Signal, iMessage, Teams, etc.) — one assistant everywhere you chat", ACCENT_CYAN),
        ("Conversation Continuity", "Persistent sessions with memory, compaction, and context carry-over — remembers what you discussed yesterday", ACCENT_GREEN),
        ("Extensible Skills Ecosystem", "ClawHub registry, bundled + managed + workspace skills, hot-reloading config — add capabilities without redeploying", ACCENT_PURPLE),
        ("Agent-to-Agent Communication", "sessions_send with ping-pong reply loops — agents can coordinate complex multi-step workflows", ACCENT_YELLOW),
        ("Voice & Visual", "Voice Wake, Talk Mode, Live Canvas (A2UI) — speak to your assistant, see visual outputs", ACCENT_CYAN),
        ("Enterprise Access Control", "DM pairing, allowlists, sendPolicy rules, Docker sandboxing — production-grade security for shared inboxes", ACCENT_GREEN),
        ("Local-First / Own Your Data", "Runs on your devices, data stays on your machine, no cloud dependency — privacy by default", ACCENT_PURPLE),
        ("Massive Community", "211k stars, 692 contributors, 48 releases — battle-tested, actively maintained, huge ecosystem", ACCENT_YELLOW),
    ]

    y = Inches(1.1)
    for i, (title, desc, color) in enumerate(advantages_oc):
        col_offset = Inches(0.4) if i % 2 == 0 else Inches(6.6)
        card = add_shape(slide, col_offset, y, Inches(6.0), Inches(1.3), BG_CARD, color)
        add_text_box(slide, col_offset + Inches(0.2), y + Inches(0.05), Inches(5.6), Inches(0.35),
                     title, font_size=14, color=color, bold=True)
        add_text_box(slide, col_offset + Inches(0.2), y + Inches(0.4), Inches(5.6), Inches(0.85),
                     desc, font_size=11, color=LIGHT_GRAY)
        if i % 2 == 1:
            y += Inches(1.45)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 8 — Advantages of SRE Ephemeral Agent
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), SRE_CLR)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "⚡ Advantages of SRE Ephemeral Agent", font_size=28, color=SRE_CLR, bold=True)

    advantages_sre = [
        ("Zero State Leakage", "Every agent starts fresh — no leftover data from previous queries. Critical for SRE where cross-contamination causes misdiagnosis", ACCENT_GREEN),
        ("Resource Efficiency", "Agents only consume memory/CPU while actively working. Zero cost when idle — perfect for cloud deployment and cost optimization", ACCENT_CYAN),
        ("Typed Specialist Agents", "7 purpose-built agent types (Log, Health, Monitor, Runbook, Trace, Dashboard, Deploy) — each knows its domain deeply", ACCENT_PURPLE),
        ("Full Lifecycle Audit Trail", "AgentRegistry tracks every agent: object ID, memory address, PID, thread, timestamps — provable creation and destruction", ACCENT_YELLOW),
        ("Natural Horizontal Scaling", "Need 10 log agents during an outage? Spawn 10 — each isolated. No shared mutable state, no locking needed", ACCENT_GREEN),
        ("Fault Isolation by Design", "If a monitoring agent crashes, log and runbook agents are unaffected. Each failure is completely contained", ACCENT_CYAN),
        ("Cooldown Pool = Best of Both", "120s cooldown means rapid-fire queries reuse the same agent (fast) but it still gets destroyed after idle (clean)", ACCENT_PURPLE),
        ("Simple Crash Recovery", "Orchestrator is stateless — just restart. No sessions to recover, no JSONL to replay, no state to reconcile", ACCENT_YELLOW),
    ]

    y = Inches(1.1)
    for i, (title, desc, color) in enumerate(advantages_sre):
        col_offset = Inches(0.4) if i % 2 == 0 else Inches(6.6)
        card = add_shape(slide, col_offset, y, Inches(6.0), Inches(1.3), BG_CARD, color)
        add_text_box(slide, col_offset + Inches(0.2), y + Inches(0.05), Inches(5.6), Inches(0.35),
                     title, font_size=14, color=color, bold=True)
        add_text_box(slide, col_offset + Inches(0.2), y + Inches(0.4), Inches(5.6), Inches(0.85),
                     desc, font_size=11, color=LIGHT_GRAY)
        if i % 2 == 1:
            y += Inches(1.45)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 9 — When to Use Which
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_YELLOW)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "When to Use Which?", font_size=28, color=WHITE, bold=True)

    # OpenClaw use cases
    add_shape(slide, Inches(0.4), Inches(1.2), Inches(5.9), Inches(5.5), BG_CARD, OPENCLAW_CLR)
    add_text_box(slide, Inches(0.7), Inches(1.3), Inches(5.3), Inches(0.5),
                 "🦞 Use OpenClaw When...", font_size=18, color=OPENCLAW_CLR, bold=True)

    oc_when = [
        ("✓ You need a personal AI assistant", 13, WHITE, True),
        ("  Always available across all your messaging apps", 11, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("✓ Conversation context matters", 13, WHITE, True),
        ("  Long-running chats that need memory of prior context", 11, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("✓ Multi-channel presence is required", 13, WHITE, True),
        ("  WhatsApp + Telegram + Slack + Discord simultaneously", 11, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("✓ Voice / visual interaction needed", 13, WHITE, True),
        ("  Talk Mode, Voice Wake, Canvas for visual outputs", 11, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("✓ Local-first privacy is critical", 13, WHITE, True),
        ("  Data never leaves your device, no cloud dependency", 11, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("✓ Complex multi-agent workflows", 13, WHITE, True),
        ("  Agents need to talk to each other and coordinate", 11, LIGHT_GRAY, False),
    ]
    add_multi_text(slide, Inches(0.7), Inches(1.9), Inches(5.3), Inches(4.8), oc_when)

    # SRE use cases
    add_shape(slide, Inches(6.8), Inches(1.2), Inches(5.9), Inches(5.5), BG_CARD, SRE_CLR)
    add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.3), Inches(0.5),
                 "⚡ Use SRE Ephemeral When...", font_size=18, color=SRE_CLR, bold=True)

    sre_when = [
        ("✓ Domain-specific operations (SRE/DevOps)", 13, WHITE, True),
        ("  Specialized agents for logs, health, monitoring, etc.", 11, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("✓ Each query must be independent", 13, WHITE, True),
        ("  SRE incidents must never cross-contaminate analysis", 11, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("✓ Cloud cost optimization matters", 13, WHITE, True),
        ("  Pay only for active work, zero idle resource usage", 11, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("✓ Compliance requires lifecycle proof", 13, WHITE, True),
        ("  Audit trail showing exactly when agents existed/worked", 11, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("✓ Parallel agent execution needed", 13, WHITE, True),
        ("  Multiple specialists working simultaneously during outage", 11, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("✓ Simple, stateless deployment", 13, WHITE, True),
        ("  Docker container with no persistent state to manage", 11, LIGHT_GRAY, False),
    ]
    add_multi_text(slide, Inches(7.1), Inches(1.9), Inches(5.3), Inches(4.8), sre_when)

    # ════════════════════════════════════════════════════════════════
    # SLIDE 10 — Key Insight / Summary
    # ════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GREEN)

    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6),
                 "Key Insight: Different Problems, Different Solutions",
                 font_size=28, color=WHITE, bold=True)

    # Analogy box
    add_shape(slide, Inches(1.5), Inches(1.4), Inches(10), Inches(1.6), BG_CARD, ACCENT_YELLOW)
    analogy = [
        ("🦞 OpenClaw = A butler who lives in your house 24/7", 18, OPENCLAW_CLR, True),
        ("  Knows your preferences, manages your life, always available", 13, LIGHT_GRAY, False),
        ("", 8, LIGHT_GRAY, False),
        ("⚡ SRE Agent = Calling a specialist when something breaks", 18, SRE_CLR, True),
        ("  Plumber shows up, fixes the pipe, sends the bill, leaves", 13, LIGHT_GRAY, False),
    ]
    add_multi_text(slide, Inches(1.8), Inches(1.5), Inches(9.4), Inches(1.5), analogy)

    # Summary points
    summary = [
        ("Both architectures are valid — they solve different problems", 16, WHITE, True),
        ("", 8, LIGHT_GRAY, False),
        ("OpenClaw excels at: conversational AI, multi-channel presence,", 14, OPENCLAW_CLR, False),
        ("  personal assistant workflows, local-first privacy", 14, OPENCLAW_CLR, False),
        ("", 8, LIGHT_GRAY, False),
        ("SRE Ephemeral excels at: operational tooling, incident response,", 14, SRE_CLR, False),
        ("  stateless cloud ops, audit compliance, resource efficiency", 14, SRE_CLR, False),
        ("", 8, LIGHT_GRAY, False),
        ("The real power? Combine both patterns:", 16, ACCENT_YELLOW, True),
        ("  Use OpenClaw as the always-on assistant layer", 14, LIGHT_GRAY, False),
        ("  Use ephemeral agents as specialized workers it can dispatch", 14, LIGHT_GRAY, False),
        ("  → sessions_spawn + typed agents = best of both worlds", 14, ACCENT_GREEN, True),
    ]
    add_multi_text(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(4), summary)

    add_text_box(slide, Inches(3), Inches(6.8), Inches(7), Inches(0.5),
                 "\"Not either/or — it's about choosing the right tool for the right job.\"",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────
    output_path = "OpenClaw_vs_SRE_Agent_Comparison.pptx"
    prs.save(output_path)
    print(f"\n✅ PPT saved: {output_path}")
    print(f"   Slides: {len(prs.slides)}")
    print("   Contents:")
    print("     1. Title Slide")
    print("     2. What Are We Comparing?")
    print("     3. Architecture: How Each System Works")
    print("     4. Core Architecture Comparison (4-Column Table)")
    print("     5. Agent Lifecycle & Session Comparison (4-Column Table)")
    print("     6. Security, Scalability & Operations (4-Column Table)")
    print("     7. 🦞 Advantages of OpenClaw")
    print("     8. ⚡ Advantages of SRE Ephemeral Agent")
    print("     9. When to Use Which?")
    print("    10. Key Insight: Summary")


if __name__ == "__main__":
    build_ppt()
